"""
File upload and text extraction endpoint.
Uploaded files are kept in memory for 1 hour then automatically purged.
Supported: PDF, DOCX, XLSX, PPTX, TXT.
"""
from __future__ import annotations

import asyncio
import io
import logging
import time
import uuid
from typing import Any

from fastapi import APIRouter, HTTPException, UploadFile, File, status
from fastapi.responses import Response

logger = logging.getLogger(__name__)
router = APIRouter(prefix="/files", tags=["files"])

# ── In-memory store: { file_id: { text, filename, expires_at } } ─────────────
_store: dict[str, dict[str, Any]] = {}
TTL_SECONDS = 3600  # 1 hour


# ── Background cleanup ────────────────────────────────────────────────────────
async def _cleanup_loop() -> None:
    while True:
        await asyncio.sleep(60)
        now = time.time()
        expired = [fid for fid, v in _store.items() if v["expires_at"] <= now]
        for fid in expired:
            _store.pop(fid, None)
        if expired:
            logger.info("Purged %d expired file(s) from memory store", len(expired))


def start_cleanup_task() -> None:
    asyncio.ensure_future(_cleanup_loop())


# ── Text extractors ───────────────────────────────────────────────────────────
def _extract_pdf(data: bytes) -> str:
    import pdfplumber  # type: ignore
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        parts = []
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                parts.append(text)
    return "\n".join(parts)


def _extract_docx(data: bytes) -> str:
    from docx import Document  # type: ignore
    doc = Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_xlsx(data: bytes) -> str:
    import openpyxl  # type: ignore
    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
    rows: list[str] = []
    for sheet in wb.worksheets:
        rows.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join("" if v is None else str(v) for v in row)
            if row_text.strip():
                rows.append(row_text)
    return "\n".join(rows)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation  # type: ignore
    prs = Presentation(io.BytesIO(data))
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _extract_txt(data: bytes) -> str:
    return data.decode("utf-8", errors="replace")


_EXTRACTORS = {
    "application/pdf": _extract_pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": _extract_docx,
    "application/msword": _extract_docx,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": _extract_xlsx,
    "application/vnd.ms-excel": _extract_xlsx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": _extract_pptx,
    "text/plain": _extract_txt,
}

_EXT_FALLBACK = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".doc": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".xls": _extract_xlsx,
    ".pptx": _extract_pptx,
    ".txt": _extract_txt,
}

MAX_FILE_BYTES = 20 * 1024 * 1024  # 20 MB


# ── Endpoint ──────────────────────────────────────────────────────────────────
@router.post("/upload", status_code=status.HTTP_201_CREATED)
async def upload_file(file: UploadFile = File(...)) -> dict:
    data = await file.read()

    if len(data) > MAX_FILE_BYTES:
        raise HTTPException(status_code=413, detail="File terlalu besar (maks 20 MB)")

    # Pick extractor
    extractor = _EXTRACTORS.get(file.content_type or "")
    if extractor is None:
        ext = ("." + file.filename.rsplit(".", 1)[-1].lower()) if file.filename and "." in file.filename else ""
        extractor = _EXT_FALLBACK.get(ext)

    if extractor is None:
        raise HTTPException(
            status_code=415,
            detail="Format file tidak didukung. Gunakan PDF, DOCX, XLSX, PPTX, atau TXT.",
        )

    try:
        text = await asyncio.get_event_loop().run_in_executor(None, extractor, data)
    except Exception as exc:
        logger.error("File extraction failed: %s", exc)
        raise HTTPException(status_code=422, detail="Gagal membaca isi file. Pastikan file tidak rusak.")

    if not text.strip():
        raise HTTPException(status_code=422, detail="File tidak mengandung teks yang dapat dibaca.")

    # Truncate to ~40k chars to stay within token limits
    if len(text) > 40_000:
        text = text[:40_000] + "\n\n[... konten dipotong karena terlalu panjang ...]"

    file_id = str(uuid.uuid4())
    _store[file_id] = {
        "file_id": file_id,
        "filename": file.filename or "file",
        "text": text,
        "expires_at": time.time() + TTL_SECONDS,
    }

    logger.info("File stored: %s (%s), %d chars", file_id, file.filename, len(text))
    return {"file_id": file_id, "filename": file.filename, "char_count": len(text)}


@router.get("/{file_id}")
async def get_file_text(file_id: str) -> dict:
    entry = _store.get(file_id)
    if not entry or entry["expires_at"] <= time.time():
        _store.pop(file_id, None)
        raise HTTPException(status_code=404, detail="File tidak ditemukan atau sudah kedaluwarsa")
    return {"file_id": file_id, "filename": entry["filename"], "text": entry["text"]}


@router.delete("/{file_id}", status_code=status.HTTP_204_NO_CONTENT, response_class=Response)
async def delete_file(file_id: str) -> Response:
    _store.pop(file_id, None)
    return Response(status_code=status.HTTP_204_NO_CONTENT)
